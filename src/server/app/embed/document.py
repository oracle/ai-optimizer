"""
Copyright (c) 2024, 2026, Oracle and/or its affiliates.
Licensed under the Universal Permissive License v1.0 as shown at http://oss.oracle.com/licenses/upl.

Document loading, splitting, and metadata processing for embedding pipelines.
"""
# spell-checker:ignore docling docos datamodel

import datetime
import functools
import json
import logging
import math
import os
from pathlib import Path
from typing import TYPE_CHECKING, Any, Optional, cast

from docling.datamodel.base_models import InputFormat
from docling.datamodel.pipeline_options import PdfPipelineOptions
from docling.document_converter import DocumentConverter, PdfFormatOption
from docling.exceptions import ConversionError
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from server.app.embed.schemas import DoclingDocumentChunk

if TYPE_CHECKING:
    from pptx.shapes.autoshape import Shape

LOGGER = logging.getLogger(__name__)


@functools.lru_cache(maxsize=2)
def _get_docling_converter(deep: bool = False) -> DocumentConverter:
    """Return a cached DocumentConverter. *deep* enables OCR and table structure."""
    pdf_pipeline_options = PdfPipelineOptions(
        do_ocr=deep,
        do_table_structure=deep,
    )

    converter = DocumentConverter(
        format_options={
            InputFormat.PDF: PdfFormatOption(pipeline_options=pdf_pipeline_options),
        }
    )
    LOGGER.info(
        "Initialized Docling DocumentConverter (deep=%s, %s)", deep, pdf_pipeline_options.model_dump(mode="json")
    )
    return converter


def _make_json_safe(value: Any) -> Any:
    """Recursively convert metadata values into JSON-serializable primitives."""
    if value is None:
        return None
    if isinstance(value, (str, int, float, bool)):
        return value
    if isinstance(value, Path):
        return str(value)
    if isinstance(value, dict):
        return {str(k): _make_json_safe(v) for k, v in value.items()}
    if isinstance(value, (list, tuple, set)):
        return [_make_json_safe(v) for v in value]
    return str(value)


def split_document(
    model: str,
    chunk_size: int,
    chunk_overlap: int,
    document: list[DoclingDocumentChunk],
) -> list[DoclingDocumentChunk]:
    """Split documents into chunks of ``chunk_size`` characters.

    ``model`` is kept for signature compatibility and logging parity.
    """
    LOGGER.info("Splitting for %s", model)
    chunk_overlap_ceil = int(math.ceil(chunk_overlap))
    if chunk_size <= 0:
        raise ValueError("chunk_size must be > 0")
    if chunk_overlap_ceil < 0:
        raise ValueError("chunk_overlap must be >= 0")
    if chunk_overlap_ceil >= chunk_size:
        raise ValueError("chunk_overlap must be smaller than chunk_size")

    def _split_text(text: str) -> list[tuple[str, int]]:
        chunks: list[tuple[str, int]] = []
        start = 0
        n = len(text)
        while start < n:
            end = min(n, start + chunk_size)
            chunk = text[start:end]
            chunks.append((chunk, start))
            if end == n:
                break
            start = end - chunk_overlap_ceil
        return chunks

    doc_split: list[DoclingDocumentChunk] = []
    for doc in document:
        text = doc.page_content or ""
        for chunk_text, start_idx in _split_text(text):
            md = dict(doc.metadata) if isinstance(doc.metadata, dict) else {}
            md["start_index"] = start_idx
            doc_split.append(DoclingDocumentChunk(page_content=chunk_text, metadata=md))

    LOGGER.info("Number of Chunks: %i", len(doc_split))
    return doc_split


def process_metadata(
    idx: int,
    chunk: DoclingDocumentChunk,
    file_metadata: Optional[dict] = None,
) -> list[DoclingDocumentChunk]:
    """Add metadata to a split document chunk."""
    filename = os.path.basename(str(chunk.metadata["source"]))

    chunk_metadata = _make_json_safe(chunk.metadata.copy())
    chunk_metadata["id"] = f"{filename}_{idx}"
    chunk_metadata["filename"] = filename

    if file_metadata and filename in file_metadata:
        chunk_metadata["size"] = file_metadata[filename].get("size")
        chunk_metadata["time_modified"] = file_metadata[filename].get("time_modified")
        chunk_metadata["etag"] = file_metadata[filename].get("etag")

    return [DoclingDocumentChunk(page_content=str(chunk.page_content), metadata=chunk_metadata)]


def _load_pdf_fast(source: str) -> list[DoclingDocumentChunk]:
    """Load a PDF using pypdf for fast text extraction."""
    reader = PdfReader(source)
    pages = [page.extract_text() or "" for page in reader.pages]
    content = "\n".join(pages)
    if not content.strip():
        raise ValueError("pypdf extracted no text content")
    return [DoclingDocumentChunk(page_content=content, metadata={"source": str(source)})]


def _load_docx_fast(source: str) -> list[DoclingDocumentChunk]:
    """Load a DOCX file using python-docx for fast text extraction."""

    doc = DocxDocument(source)
    content = "\n".join(para.text for para in doc.paragraphs)
    if not content.strip():
        raise ValueError("python-docx extracted no text content")
    return [DoclingDocumentChunk(page_content=content, metadata={"source": str(source)})]


def _load_pptx_fast(source: str) -> list[DoclingDocumentChunk]:
    """Load a PPTX file using python-pptx for fast text extraction."""

    prs = Presentation(source)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(cast("Shape", shape).text_frame.text)
    content = "\n".join(texts)
    if not content.strip():
        raise ValueError("python-pptx extracted no text content")
    return [DoclingDocumentChunk(page_content=content, metadata={"source": str(source)})]


def _load_xlsx_fast(source: str) -> list[DoclingDocumentChunk]:
    """Load an XLSX file using openpyxl for fast text extraction."""

    wb = load_workbook(source, read_only=True, data_only=True)
    rows = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            rows.append("\t".join(str(c) if c is not None else "" for c in row))
    wb.close()
    content = "\n".join(rows)
    if not content.strip():
        raise ValueError("openpyxl extracted no text content")
    return [DoclingDocumentChunk(page_content=content, metadata={"source": str(source)})]


# Map file extensions to their fast loader functions
_FAST_LOADERS: dict[str, Any] = {
    ".pdf": _load_pdf_fast,
    ".docx": _load_docx_fast,
    ".pptx": _load_pptx_fast,
    ".xlsx": _load_xlsx_fast,
}


def _load_docling_document(source: str, parsing_mode: str = "fast") -> list[DoclingDocumentChunk]:
    """Load a file (or URL) and return a single-document list.

    In *fast* mode, tries lightweight library extraction first (pypdf, python-docx,
    python-pptx, openpyxl) and falls back to Docling on failure.
    In *deep* mode, always uses Docling with OCR and table-structure analysis enabled.
    """
    source_str = str(source)
    is_remote = source_str.startswith(("http://", "https://"))
    ext = os.path.splitext(source_str)[1].lower()

    if parsing_mode == "fast" and not is_remote:
        fast_loader = _FAST_LOADERS.get(ext)
        if fast_loader is not None:
            try:
                result = fast_loader(source)
                LOGGER.info("%s loaded via fast path (%s)", ext, fast_loader.__name__)
                return result
            except Exception as exc:
                LOGGER.warning("Fast loader failed for %s (%s), falling back to Docling", source, exc)

    deep = parsing_mode == "deep"
    converter = _get_docling_converter(deep=deep)
    result = converter.convert(source=source)
    dl_doc = result.document
    try:
        content = dl_doc.export_to_markdown()
    except Exception:
        content = dl_doc.export_to_text()

    return [DoclingDocumentChunk(page_content=str(content), metadata={"source": str(source)})]


def _capture_file_metadata(name: str, stat: os.stat_result, file_metadata: dict) -> None:
    """Capture file metadata if not already provided."""
    if name not in file_metadata:
        file_metadata[name] = {
            "size": stat.st_size,
            "time_modified": datetime.datetime.fromtimestamp(stat.st_mtime, datetime.timezone.utc).isoformat(),
        }


def _process_and_split_document(
    loaded_doc: list[DoclingDocumentChunk],
    split: bool,
    model: str,
    chunk_size: int,
    chunk_overlap: int,
    file_metadata: dict,
) -> list[DoclingDocumentChunk]:
    """Process and split a loaded document."""
    if not split:
        return loaded_doc

    split_doc = split_document(model, chunk_size, chunk_overlap, loaded_doc)
    split_docos: list[DoclingDocumentChunk] = []
    for idx, chunk in enumerate(split_doc, start=1):
        split_doc_with_mdata = process_metadata(idx, chunk, file_metadata)
        split_docos += split_doc_with_mdata
    return split_docos


def load_and_split_documents(
    src_files: list,
    model: str,
    chunk_size: int,
    chunk_overlap: int,
    write_json: bool = False,
    output_dir: Optional[str] = None,
    file_metadata: Optional[dict] = None,
    parsing_mode: str = "fast",
) -> tuple[list[DoclingDocumentChunk], list, dict]:
    """Load files via Docling, split into chunks, and return processing results.

    Returns:
        Tuple of (chunks, json_file_paths, processing_results_dict).
    """
    split_files: list = []
    all_split_docos: list[DoclingDocumentChunk] = []
    processing_results: dict = {"processed_files": [], "skipped_files": [], "total_chunks": 0}

    if file_metadata is None:
        file_metadata = {}

    for file in src_files:
        stat = os.stat(file)
        LOGGER.info("Loading %s (%i bytes)", os.path.basename(file), stat.st_size)

        _capture_file_metadata(os.path.basename(file), stat, file_metadata)

        try:
            loaded_doc = _load_docling_document(file, parsing_mode=parsing_mode)
            LOGGER.info("Loaded Pages: %i", len(loaded_doc))

            split_docos = _process_and_split_document(loaded_doc, True, model, chunk_size, chunk_overlap, file_metadata)

            if write_json and output_dir:
                split_files.append(doc_to_json(split_docos, file, output_dir))

            all_split_docos += split_docos
            processing_results["processed_files"].append(
                {"filename": os.path.basename(file), "chunks": len(split_docos)}
            )

        except (ValueError, ConversionError) as e:
            LOGGER.warning("Skipping unsupported file %s: %s", os.path.basename(file), str(e))
            processing_results["skipped_files"].append(
                {
                    "filename": os.path.basename(file),
                    "reason": f"Unsupported file type: {os.path.splitext(file)[1][1:]}",
                }
            )
            continue
        except Exception as e:
            LOGGER.warning("Skipping file %s due to processing error: %s", os.path.basename(file), str(e))
            processing_results["skipped_files"].append(
                {"filename": os.path.basename(file), "reason": f"Processing error: {str(e)}"}
            )
            continue

    processing_results["total_chunks"] = len(all_split_docos)
    LOGGER.info("Total Number of Chunks: %i", len(all_split_docos))
    LOGGER.info(
        "Processed files: %i, Skipped files: %i",
        len(processing_results["processed_files"]),
        len(processing_results["skipped_files"]),
    )

    return all_split_docos, split_files, processing_results


def doc_to_json(document: list[DoclingDocumentChunk], file: str, output_dir: str) -> str:
    """Write document chunks to a JSON file. Returns the destination file path."""
    src_file_name = os.path.basename(file)
    dst_file_name = "_" + os.path.splitext(src_file_name)[0] + ".json"

    docs_dict = [doc.to_json() for doc in document]
    json_data = json.dumps(docs_dict, indent=4)

    dst_file_path = os.path.join(output_dir, dst_file_name)
    with open(dst_file_path, "w", encoding="utf-8") as f:
        f.write(json_data)
    file_size = os.path.getsize(dst_file_path)
    LOGGER.info("Wrote split JSON file: %s (%i bytes)", dst_file_path, file_size)

    return dst_file_path


def json_to_doc(file: str) -> list[DoclingDocumentChunk]:
    """Create a list of DoclingDocumentChunk from a JSON file."""
    LOGGER.info("Converting %s to Document", file)

    with open(file, "r", encoding="utf-8") as document:
        chunks = json.load(document)
        docs = []
        for chunk in chunks:
            page_content = chunk["kwargs"]["page_content"]
            metadata = chunk["kwargs"]["metadata"]
            docs.append(DoclingDocumentChunk(page_content=str(page_content), metadata=metadata))

    LOGGER.info("Chunks ingested: %i", len(docs))
    return docs
