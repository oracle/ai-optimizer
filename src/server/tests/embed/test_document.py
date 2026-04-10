"""
Copyright (c) 2024, 2026, Oracle and/or its affiliates.
Licensed under the Universal Permissive License v1.0 as shown at http://oss.oracle.com/licenses/upl.

Tests for embed document processing utilities.
"""
# spell-checker: disable

from unittest.mock import MagicMock, patch

import pytest

from server.app.embed.document import (
    DoclingDocumentChunk,
    _load_docling_document,
    _load_docx_fast,
    _load_pdf_fast,
    _load_pptx_fast,
    _load_xlsx_fast,
    _make_json_safe,
    load_and_split_documents,
    process_metadata,
    split_document,
)

MODULE = "server.app.embed.document"


# ---------------------------------------------------------------------------
# DoclingDocumentChunk
# ---------------------------------------------------------------------------


@pytest.mark.unit
def test_chunk_creation():
    """DoclingDocumentChunk creates with expected fields."""
    chunk = DoclingDocumentChunk(page_content="hello", metadata={"source": "test.pdf"})
    assert chunk.page_content == "hello"
    assert chunk.metadata["source"] == "test.pdf"
    assert chunk.id is None


@pytest.mark.unit
def test_chunk_id_from_metadata():
    """DoclingDocumentChunk extracts id from metadata."""
    chunk = DoclingDocumentChunk(page_content="hello", metadata={"id": "abc_1"})
    assert chunk.id == "abc_1"


@pytest.mark.unit
def test_chunk_explicit_id():
    """DoclingDocumentChunk uses explicit doc_id over metadata."""
    chunk = DoclingDocumentChunk(page_content="hello", metadata={"id": "meta_id"}, doc_id="explicit_id")
    assert chunk.id == "explicit_id"


@pytest.mark.unit
def test_chunk_to_json():
    """to_json produces LangChain-compatible shape."""
    chunk = DoclingDocumentChunk(page_content="content", metadata={"key": "val"})
    result = chunk.to_json()
    assert result["lc"] == 1
    assert result["type"] == "constructor"
    assert result["kwargs"]["page_content"] == "content"
    assert result["kwargs"]["metadata"]["key"] == "val"


# ---------------------------------------------------------------------------
# _make_json_safe
# ---------------------------------------------------------------------------


@pytest.mark.unit
def test_make_json_safe_primitives():
    """Primitives pass through unchanged."""
    assert _make_json_safe("text") == "text"
    assert _make_json_safe(42) == 42
    assert _make_json_safe(True) is True
    assert _make_json_safe(None) is None


@pytest.mark.unit
def test_make_json_safe_nested():
    """Nested structures are recursively converted."""
    from pathlib import Path

    data = {"path": Path("/tmp/test"), "items": [1, Path("/a")]}
    result = _make_json_safe(data)
    assert result["path"] == "/tmp/test"
    assert result["items"] == [1, "/a"]


# ---------------------------------------------------------------------------
# split_document
# ---------------------------------------------------------------------------


@pytest.mark.unit
def test_split_document_basic():
    """Splits document into expected number of chunks."""
    doc = DoclingDocumentChunk(page_content="a" * 100, metadata={"source": "test"})
    result = split_document("test-model", 30, 10, [doc])
    # 100 chars, chunk_size=30, overlap=10 → chunks at 0-30, 20-50, 40-70, 60-90, 80-100
    assert len(result) == 5
    assert result[0].metadata["start_index"] == 0


@pytest.mark.unit
def test_split_document_preserves_metadata():
    """Split chunks inherit source metadata."""
    doc = DoclingDocumentChunk(page_content="Hello world content", metadata={"source": "doc.pdf", "page": 1})
    result = split_document("test-model", 100, 0, [doc])
    assert len(result) == 1
    assert result[0].metadata["source"] == "doc.pdf"
    assert result[0].metadata["page"] == 1


@pytest.mark.unit
def test_split_document_invalid_params():
    """Raises ValueError for invalid chunk parameters."""
    doc = DoclingDocumentChunk(page_content="text", metadata={})
    with pytest.raises(ValueError, match="chunk_size must be > 0"):
        split_document("model", 0, 0, [doc])
    with pytest.raises(ValueError, match="chunk_overlap must be smaller"):
        split_document("model", 10, 10, [doc])


# ---------------------------------------------------------------------------
# process_metadata
# ---------------------------------------------------------------------------


@pytest.mark.unit
def test_process_metadata_basic():
    """Adds filename and id to metadata."""
    chunk = DoclingDocumentChunk(page_content="text", metadata={"source": "/path/to/test.pdf", "start_index": 0})
    result = process_metadata(1, chunk)
    assert len(result) == 1
    assert result[0].metadata["filename"] == "test.pdf"
    assert result[0].metadata["id"] == "test.pdf_1"


@pytest.mark.unit
def test_process_metadata_with_file_metadata():
    """Enriches chunk with file size, time, and etag."""
    chunk = DoclingDocumentChunk(page_content="text", metadata={"source": "doc.pdf"})
    file_meta = {"doc.pdf": {"size": 1024, "time_modified": "2026-01-01T00:00:00", "etag": "abc"}}
    result = process_metadata(1, chunk, file_meta)
    assert result[0].metadata["size"] == 1024
    assert result[0].metadata["etag"] == "abc"


# ---------------------------------------------------------------------------
# Fast loaders — real file tests
# ---------------------------------------------------------------------------


def _create_minimal_pdf(path: str, text: str = "Hello PDF") -> str:
    """Create a minimal PDF with text content using pypdf."""
    from pypdf import PdfWriter
    from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    page = writer.pages[0]

    stream = DecodedStreamObject()
    stream.set_data(f"BT /F1 12 Tf 10 50 Td ({text}) Tj ET".encode())

    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_ref = writer._add_object(font)
    resources = DictionaryObject({NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_ref})})
    dict.__setitem__(page, NameObject("/Resources"), resources)
    dict.__setitem__(page, NameObject("/Contents"), writer._add_object(stream))

    with open(path, "wb") as f:
        writer.write(f)
    return path


def _create_minimal_docx(path: str, text: str = "Hello DOCX") -> str:
    """Create a minimal .docx with a single paragraph."""
    from docx import Document as DocxDocument

    doc = DocxDocument()
    doc.add_paragraph(text)
    doc.save(path)
    return path


def _create_minimal_pptx(path: str, text: str = "Hello PPTX") -> str:
    """Create a minimal .pptx with a single text box."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    text_box.text_frame.text = text
    prs.save(path)
    return path


def _create_minimal_xlsx(path: str, text: str = "Hello XLSX") -> str:
    """Create a minimal .xlsx with a single cell."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws["A1"] = text
    wb.save(path)
    wb.close()
    return path


@pytest.mark.unit
def test_load_pdf_fast(tmp_path):
    """_load_pdf_fast extracts text from a real PDF."""
    pdf_path = str(tmp_path / "test.pdf")
    _create_minimal_pdf(pdf_path, "Hello PDF")
    result = _load_pdf_fast(pdf_path)
    assert len(result) == 1
    assert "Hello PDF" in result[0].page_content
    assert result[0].metadata["source"] == pdf_path


@pytest.mark.unit
def test_load_pdf_fast_empty(tmp_path):
    """_load_pdf_fast raises ValueError on empty PDF."""
    from pypdf import PdfWriter

    pdf_path = str(tmp_path / "empty.pdf")
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    with open(pdf_path, "wb") as f:
        writer.write(f)
    with pytest.raises(ValueError, match="no text content"):
        _load_pdf_fast(pdf_path)


@pytest.mark.unit
def test_load_docx_fast(tmp_path):
    """_load_docx_fast extracts text from a real DOCX."""
    docx_path = str(tmp_path / "test.docx")
    _create_minimal_docx(docx_path, "Hello DOCX")
    result = _load_docx_fast(docx_path)
    assert len(result) == 1
    assert "Hello DOCX" in result[0].page_content
    assert result[0].metadata["source"] == docx_path


@pytest.mark.unit
def test_load_docx_fast_empty(tmp_path):
    """_load_docx_fast raises ValueError on empty DOCX."""
    from docx import Document as DocxDocument

    docx_path = str(tmp_path / "empty.docx")
    doc = DocxDocument()
    doc.save(docx_path)
    with pytest.raises(ValueError, match="no text content"):
        _load_docx_fast(docx_path)


@pytest.mark.unit
def test_load_pptx_fast(tmp_path):
    """_load_pptx_fast extracts text from a real PPTX."""
    pptx_path = str(tmp_path / "test.pptx")
    _create_minimal_pptx(pptx_path, "Hello PPTX")
    result = _load_pptx_fast(pptx_path)
    assert len(result) == 1
    assert "Hello PPTX" in result[0].page_content
    assert result[0].metadata["source"] == pptx_path


@pytest.mark.unit
def test_load_pptx_fast_empty(tmp_path):
    """_load_pptx_fast raises ValueError on empty PPTX."""
    from pptx import Presentation

    pptx_path = str(tmp_path / "empty.pptx")
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(pptx_path)
    with pytest.raises(ValueError, match="no text content"):
        _load_pptx_fast(pptx_path)


@pytest.mark.unit
def test_load_xlsx_fast(tmp_path):
    """_load_xlsx_fast extracts text from a real XLSX."""
    xlsx_path = str(tmp_path / "test.xlsx")
    _create_minimal_xlsx(xlsx_path, "Hello XLSX")
    result = _load_xlsx_fast(xlsx_path)
    assert len(result) == 1
    assert "Hello XLSX" in result[0].page_content
    assert result[0].metadata["source"] == xlsx_path


@pytest.mark.unit
def test_load_xlsx_fast_empty(tmp_path):
    """_load_xlsx_fast raises ValueError on empty XLSX."""
    from openpyxl import Workbook

    xlsx_path = str(tmp_path / "empty.xlsx")
    wb = Workbook()
    wb.save(xlsx_path)
    wb.close()
    with pytest.raises(ValueError, match="no text content"):
        _load_xlsx_fast(xlsx_path)


# ---------------------------------------------------------------------------
# _load_docling_document — dispatch logic
# ---------------------------------------------------------------------------


@pytest.mark.unit
def test_dispatch_fast_mode_uses_fast_loader(tmp_path):
    """Fast mode with a known extension calls the fast loader, not Docling."""
    pdf_path = str(tmp_path / "test.pdf")
    _create_minimal_pdf(pdf_path, "fast path")
    with patch(f"{MODULE}._get_docling_converter") as mock_converter:
        result = _load_docling_document(pdf_path, parsing_mode="fast")
    mock_converter.assert_not_called()
    assert "fast path" in result[0].page_content


@pytest.mark.unit
def test_dispatch_fast_mode_fallback_on_failure(tmp_path):
    """Fast mode falls back to Docling when fast loader raises."""
    pdf_path = str(tmp_path / "test.pdf")
    _create_minimal_pdf(pdf_path, "content")

    mock_doc = MagicMock()
    mock_doc.export_to_markdown.return_value = "docling content"
    mock_result = MagicMock()
    mock_result.document = mock_doc
    mock_converter = MagicMock()
    mock_converter.convert.return_value = mock_result

    broken_loader = MagicMock(side_effect=RuntimeError("broken"))
    patched_loaders = {".pdf": broken_loader}

    with (
        patch(f"{MODULE}._FAST_LOADERS", patched_loaders),
        patch(f"{MODULE}._get_docling_converter", return_value=mock_converter),
    ):
        result = _load_docling_document(pdf_path, parsing_mode="fast")
    broken_loader.assert_called_once()
    assert result[0].page_content == "docling content"


@pytest.mark.unit
def test_dispatch_fast_mode_unknown_ext_uses_docling(tmp_path):
    """Fast mode with no fast loader (e.g. .html) goes straight to Docling."""
    html_path = str(tmp_path / "page.html")
    html_path_obj = tmp_path / "page.html"
    html_path_obj.write_text("<html><body>Hello</body></html>")

    mock_doc = MagicMock()
    mock_doc.export_to_markdown.return_value = "Hello"
    mock_result = MagicMock()
    mock_result.document = mock_doc
    mock_converter = MagicMock()
    mock_converter.convert.return_value = mock_result

    with patch(f"{MODULE}._get_docling_converter", return_value=mock_converter):
        result = _load_docling_document(html_path, parsing_mode="fast")
    mock_converter.convert.assert_called_once()
    assert result[0].page_content == "Hello"


@pytest.mark.unit
def test_dispatch_deep_mode_skips_fast_loader(tmp_path):
    """Deep mode always uses Docling even for extensions with fast loaders."""
    pdf_path = str(tmp_path / "test.pdf")
    _create_minimal_pdf(pdf_path, "content")

    mock_doc = MagicMock()
    mock_doc.export_to_markdown.return_value = "deep content"
    mock_result = MagicMock()
    mock_result.document = mock_doc
    mock_converter = MagicMock()
    mock_converter.convert.return_value = mock_result

    with (
        patch(f"{MODULE}._load_pdf_fast") as mock_fast,
        patch(f"{MODULE}._get_docling_converter", return_value=mock_converter) as mock_get,
    ):
        result = _load_docling_document(pdf_path, parsing_mode="deep")
    mock_fast.assert_not_called()
    mock_get.assert_called_once_with(deep=True)
    assert result[0].page_content == "deep content"


@pytest.mark.unit
@pytest.mark.parametrize("ext", [".html", ".md", ".csv", ".txt"])
def test_dispatch_deep_mode_unknown_ext(tmp_path, ext):
    """Deep mode uses Docling with deep=True for non-fast-loader types."""
    file_path = str(tmp_path / f"file{ext}")
    (tmp_path / f"file{ext}").write_text("some content")

    mock_doc = MagicMock()
    mock_doc.export_to_markdown.return_value = "deep content"
    mock_result = MagicMock()
    mock_result.document = mock_doc
    mock_converter = MagicMock()
    mock_converter.convert.return_value = mock_result

    with patch(f"{MODULE}._get_docling_converter", return_value=mock_converter) as mock_get:
        _load_docling_document(file_path, parsing_mode="deep")
    mock_get.assert_called_once_with(deep=True)


# ---------------------------------------------------------------------------
# load_and_split_documents — per file-type with fast loaders
# ---------------------------------------------------------------------------


@pytest.mark.unit
@pytest.mark.parametrize(
    "ext, creator, needle",
    [
        (".pdf", _create_minimal_pdf, "Hello PDF"),
        (".docx", _create_minimal_docx, "Hello DOCX"),
        (".pptx", _create_minimal_pptx, "Hello PPTX"),
        (".xlsx", _create_minimal_xlsx, "Hello XLSX"),
    ],
    ids=["pdf", "docx", "pptx", "xlsx"],
)
def test_load_and_split_fast_loaders(tmp_path, ext, creator, needle):
    """load_and_split_documents processes each fast-loader type in fast mode."""
    file_path = str(tmp_path / f"test{ext}")
    creator(file_path, needle)
    chunks, _, results = load_and_split_documents(
        src_files=[file_path],
        model="test-model",
        chunk_size=5000,
        chunk_overlap=0,
        parsing_mode="fast",
    )
    assert len(results["processed_files"]) == 1
    assert len(results["skipped_files"]) == 0
    assert any(needle in c.page_content for c in chunks)


@pytest.mark.unit
@pytest.mark.parametrize("ext", [".html", ".md", ".csv", ".txt"])
def test_load_and_split_docling_types(tmp_path, ext):
    """load_and_split_documents processes Docling-only types via mocked converter."""
    file_path = str(tmp_path / f"test{ext}")
    (tmp_path / f"test{ext}").write_text("test content for " + ext)

    mock_doc = MagicMock()
    mock_doc.export_to_markdown.return_value = "docling parsed " + ext
    mock_result = MagicMock()
    mock_result.document = mock_doc
    mock_converter = MagicMock()
    mock_converter.convert.return_value = mock_result

    with patch(f"{MODULE}._get_docling_converter", return_value=mock_converter):
        chunks, _, results = load_and_split_documents(
            src_files=[file_path],
            model="test-model",
            chunk_size=5000,
            chunk_overlap=0,
            parsing_mode="fast",
        )
    assert len(results["processed_files"]) == 1
    assert len(results["skipped_files"]) == 0
    assert any(("docling parsed " + ext) in c.page_content for c in chunks)


@pytest.mark.unit
@pytest.mark.parametrize(
    "ext, creator, needle",
    [
        (".pdf", _create_minimal_pdf, "Hello PDF"),
        (".docx", _create_minimal_docx, "Hello DOCX"),
        (".pptx", _create_minimal_pptx, "Hello PPTX"),
        (".xlsx", _create_minimal_xlsx, "Hello XLSX"),
    ],
    ids=["pdf", "docx", "pptx", "xlsx"],
)
def test_load_and_split_deep_mode(tmp_path, ext, creator, needle):
    """Deep mode uses Docling for all types, including those with fast loaders."""
    file_path = str(tmp_path / f"test{ext}")
    creator(file_path, needle)

    mock_doc = MagicMock()
    mock_doc.export_to_markdown.return_value = "deep " + needle
    mock_result = MagicMock()
    mock_result.document = mock_doc
    mock_converter = MagicMock()
    mock_converter.convert.return_value = mock_result

    with patch(f"{MODULE}._get_docling_converter", return_value=mock_converter):
        chunks, _, results = load_and_split_documents(
            src_files=[file_path],
            model="test-model",
            chunk_size=5000,
            chunk_overlap=0,
            parsing_mode="deep",
        )
    assert len(results["processed_files"]) == 1
    assert any(("deep " + needle) in c.page_content for c in chunks)


@pytest.mark.unit
@pytest.mark.parametrize("ext", [".png", ".jpg", ".jpeg"])
def test_load_and_split_image_types(tmp_path, ext):
    """Image types are processed via Docling (which handles OCR in deep mode)."""
    # Create a minimal 1x1 PNG (images need binary content, not text)
    file_path = str(tmp_path / f"image{ext}")
    # Write a minimal valid PNG for all image types (parsers check magic bytes)
    import struct
    import zlib

    # Minimal 1x1 white PNG
    def _minimal_png():
        signature = b"\x89PNG\r\n\x1a\n"
        ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
        ihdr_crc = struct.pack(">I", zlib.crc32(b"IHDR" + ihdr_data) & 0xFFFFFFFF)
        ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + ihdr_crc
        raw = zlib.compress(b"\x00\xff\xff\xff")
        idat_crc = struct.pack(">I", zlib.crc32(b"IDAT" + raw) & 0xFFFFFFFF)
        idat = struct.pack(">I", len(raw)) + b"IDAT" + raw + idat_crc
        iend_crc = struct.pack(">I", zlib.crc32(b"IEND") & 0xFFFFFFFF)
        iend = struct.pack(">I", 0) + b"IEND" + iend_crc
        return signature + ihdr + idat + iend

    with open(file_path, "wb") as f:
        f.write(_minimal_png())

    mock_doc = MagicMock()
    mock_doc.export_to_markdown.return_value = "image text via ocr"
    mock_result = MagicMock()
    mock_result.document = mock_doc
    mock_converter = MagicMock()
    mock_converter.convert.return_value = mock_result

    with patch(f"{MODULE}._get_docling_converter", return_value=mock_converter):
        _, _, results = load_and_split_documents(
            src_files=[file_path],
            model="test-model",
            chunk_size=5000,
            chunk_overlap=0,
            parsing_mode="deep",
        )
    assert len(results["processed_files"]) == 1
